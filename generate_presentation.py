#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Kitchensink Application
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Objective
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    title_shape = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Kitchensink User Management System"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    content_shape = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    content_frame = content_shape.text_frame
    content_frame.text = "Objective"
    content_para = content_frame.paragraphs[0]
    content_para.font.size = Pt(32)
    content_para.font.bold = True
    
    p = content_frame.add_paragraph()
    p.text = "• Migrate Jakarta EE application to Spring Boot with MongoDB"
    p.font.size = Pt(20)
    p.space_after = Pt(12)
    
    p = content_frame.add_paragraph()
    p.text = "• Build modern REST API with comprehensive security features"
    p.font.size = Pt(20)
    p.space_after = Pt(12)
    
    p = content_frame.add_paragraph()
    p.text = "• Implement user management system with OTP-based authentication"
    p.font.size = Pt(20)
    p.space_after = Pt(12)
    
    p = content_frame.add_paragraph()
    p.text = "• Add enterprise-grade features: caching, rate limiting, audit logging"
    p.font.size = Pt(20)
    p.space_after = Pt(12)
    
    p = content_frame.add_paragraph()
    p.text = "• Create React frontend for seamless user experience"
    p.font.size = Pt(20)
    
    # Slide 2: Existing Features vs Added Features
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide2.shapes.title
    title.text = "Existing Features vs Added Features"
    
    content = slide2.placeholders[1]
    tf = content.text_frame
    tf.text = "Existing Features (Migrated from EJB)"
    
    p = tf.add_paragraph()
    p.text = "• Create Member (POST)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Get All Members (GET)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Get Member by ID (GET)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• View members ordered by name"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "New Features Added"
    p.space_before = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Update Member (PUT) & Delete Member (DELETE)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Search by name (fuzzy search) & Filter by email domain"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• OTP-based authentication & JWT token management"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Cursor-based pagination & Role-based access control"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Update request workflow (approve/reject)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• PII encryption, Audit logging, Rate limiting, Caching"
    p.level = 1
    
    # Slide 3: Tech Stack
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide3.shapes.title
    title.text = "Technology Stack"
    
    content = slide3.placeholders[1]
    tf = content.text_frame
    tf.text = "Backend"
    
    p = tf.add_paragraph()
    p.text = "• Spring Boot 3.5.7, Java 21, Maven"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• MongoDB 8.0+ (NoSQL Database)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Spring Security + JWT (Authentication)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Caffeine Cache, Jasypt (Encryption)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Spring Actuator (Health Checks)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Swagger/OpenAPI (Documentation)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Frontend"
    p.space_before = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• React 18, Axios, React Router DOM"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Infrastructure"
    p.space_before = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• MongoDB Collections, SMTP Email Service"
    p.level = 1
    
    # Slide 4: Caching Implementation
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide4.shapes.title
    title.text = "Caching Implementation"
    
    content = slide4.placeholders[1]
    tf = content.text_frame
    tf.text = "Caffeine Cache - In-Memory Caching"
    
    p = tf.add_paragraph()
    p.text = "• Three cache instances configured:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  - userById: 500 entries, 3min write TTL, 1min access TTL"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "  - roleById: 50 entries, 15min write TTL, 5min access TTL"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "  - roleNameByUserId: 500 entries, 3min write TTL"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• @Cacheable annotation on service methods"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Automatic cache eviction based on TTL"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Cache statistics enabled for monitoring"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Reduces database queries for frequently accessed data"
    p.level = 1
    
    # Slide 5: Health Check Implementation
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide5.shapes.title
    title.text = "Health Check Implementation"
    
    content = slide5.placeholders[1]
    tf = content.text_frame
    tf.text = "Spring Boot Actuator"
    
    p = tf.add_paragraph()
    p.text = "• Exposed endpoints: health, info, metrics, prometheus"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Health endpoint: /actuator/health"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• MongoDB health check enabled"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Detailed health info (when authorized)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Prometheus metrics for monitoring"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Application info: name, version, Java version"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Ready for integration with monitoring tools"
    p.level = 1
    
    # Slide 6: Security Implementation
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide6.shapes.title
    title.text = "Security Implementation"
    
    content = slide6.placeholders[1]
    tf = content.text_frame
    tf.text = "Multi-Layer Security"
    
    p = tf.add_paragraph()
    p.text = "JWT Authentication:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  - Stateless token-based auth, 24hr expiration"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "  - Role-based access (ADMIN, USER)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "  - Method-level security with @PreAuthorize"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "API Key Authentication:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  - X-API-Key header validation"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "  - Configurable via environment variables"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "  - Public endpoints excluded (actuator, swagger)"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "CORS Protection:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  - Whitelist-based origin validation"
    p.level = 2
    
    # Slide 7: Rate Limiter Implementation
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide7.shapes.title
    title.text = "Rate Limiter Implementation"
    
    content = slide7.placeholders[1]
    tf = content.text_frame
    tf.text = "Request Throttling"
    
    p = tf.add_paragraph()
    p.text = "• RateLimitFilter - IP-based rate limiting"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Default: 60 requests per minute per IP"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Sliding window algorithm with ConcurrentHashMap"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Automatic window reset after 1 minute"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Returns HTTP 429 (Too Many Requests) when exceeded"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Excludes: actuator, swagger, api-docs endpoints"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "OTP Rate Limiting:"
    p.space_before = Pt(12)
    p = tf.add_paragraph()
    p.text = "• 1000 attempts per 15 minutes per email"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Prevents OTP abuse and brute force attacks"
    p.level = 1
    
    # Slide 8: Pagination Implementation
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide8.shapes.title
    title.text = "Pagination Implementation"
    
    content = slide8.placeholders[1]
    tf = content.text_frame
    tf.text = "Cursor-Based Pagination"
    
    p = tf.add_paragraph()
    p.text = "• CursorPageResponse with next/previous cursors"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Supports forward (next) and backward (previous) navigation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Sort by ID (default) or name"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Fetches size+1 records to determine hasNext/hasPrevious"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Efficient for large datasets (no offset calculation)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Backward compatible with page-based pagination"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Max page size: 100 records"
    p.level = 1
    
    # Slide 9: Audit Logging Implementation
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide9.shapes.title
    title.text = "Audit Logging Implementation"
    
    content = slide9.placeholders[1]
    tf = content.text_frame
    tf.text = "Comprehensive Change Tracking"
    
    p = tf.add_paragraph()
    p.text = "• UserMongoEventListener - Automatic audit on User changes"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Tracks CREATE, UPDATE, DELETE operations"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Field-level change tracking (old vs new values)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Captures metadata: timestamp, IP address, correlation ID"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Stores in audit_logs MongoDB collection"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• ThreadLocal for tracking old state before updates"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Non-blocking async logging for performance"
    p.level = 1
    
    # Slide 10: PII Encryption Implementation
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide10.shapes.title
    title.text = "PII Encryption Implementation"
    
    content = slide10.placeholders[1]
    tf = content.text_frame
    tf.text = "Jasypt Encryption Service"
    
    p = tf.add_paragraph()
    p.text = "• PBEWITHHMACSHA512ANDAES_256 algorithm"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Encrypts: email, phone numbers (PII data)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Key versioning support (v1:encryptedData format)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Legacy key support for key rotation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Automatic decryption on read operations"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Random salt and IV generation for each encryption"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Configurable via environment variables"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• OTP values hashed (SHA-256) before storage"
    p.level = 1
    
    # Slide 11: Input Sanitization Implementation
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide11.shapes.title
    title.text = "Input Sanitization Implementation"
    
    content = slide11.placeholders[1]
    tf = content.text_frame
    tf.text = "XSS Prevention"
    
    p = tf.add_paragraph()
    p.text = "• InputSanitizationService using Apache Commons Text"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Escapes HTML/XML special characters"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Sanitizes user input before database storage"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Prevents script injection attacks"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Applied to all user-provided text fields"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Works in conjunction with Spring Validation"
    p.level = 1
    
    # Slide 12: Email & OTP Service Implementation
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide12.shapes.title
    title.text = "Email & OTP Service Implementation"
    
    content = slide12.placeholders[1]
    tf = content.text_frame
    tf.text = "OTP-Based Authentication"
    
    p = tf.add_paragraph()
    p.text = "OTP Service:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  - 6-digit numeric OTP generation"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "  - 10-minute expiration"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "  - SHA-256 hashing before storage"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "  - Rate limiting: 1000 attempts/15min"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Email Service:"
    p.space_before = Pt(12)
    p = tf.add_paragraph()
    p.text = "  - Spring Mail with SMTP (Gmail)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  - Sends OTP codes for login/email change"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  - Update request notifications (admin/user)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  - Async email processing"
    p.level = 1
    
    # Slide 13: Architecture Diagram
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    title_shape = slide13.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Backend Architecture"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    content_shape = slide13.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
    content_frame = content_shape.text_frame
    content_frame.text = "Request Flow (Detailed diagram available in archDiag5.pdf):"
    content_para = content_frame.paragraphs[0]
    content_para.font.size = Pt(20)
    content_para.font.bold = True
    
    p = content_frame.add_paragraph()
    p.text = "1. HTTP Request → CORS Filter → Correlation ID Filter"
    p.font.size = Pt(15)
    p.space_after = Pt(5)
    
    p = content_frame.add_paragraph()
    p.text = "2. Rate Limit Filter → Request Logging Filter"
    p.font.size = Pt(15)
    p.space_after = Pt(5)
    
    p = content_frame.add_paragraph()
    p.text = "3. JWT/API Key Authentication Filter"
    p.font.size = Pt(15)
    p.space_after = Pt(5)
    
    p = content_frame.add_paragraph()
    p.text = "4. Controllers (Auth/Profile/Admin) → Service Layer"
    p.font.size = Pt(15)
    p.space_after = Pt(5)
    
    p = content_frame.add_paragraph()
    p.text = "5. Cache Check (Caffeine) → Repository Layer"
    p.font.size = Pt(15)
    p.space_after = Pt(5)
    
    p = content_frame.add_paragraph()
    p.text = "6. MongoDB (users, roles, otps, audit_logs collections)"
    p.font.size = Pt(15)
    p.space_after = Pt(5)
    
    p = content_frame.add_paragraph()
    p.text = "7. Event Listeners → Audit Service → Response"
    p.font.size = Pt(15)
    p.space_after = Pt(8)
    
    p = content_frame.add_paragraph()
    p.text = "Layers: Security → Controller → Service → Repository → Database"
    p.font.size = Pt(13)
    p.font.italic = True
    p.space_after = Pt(5)
    
    p = content_frame.add_paragraph()
    p.text = "Key Features: PII Encryption, Input Sanitization, Email Service, Async Processing"
    p.font.size = Pt(13)
    p.font.italic = True
    
    # Slide 14: Scaling the System
    slide14 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide14.shapes.title
    title.text = "Scaling the System"
    
    content = slide14.placeholders[1]
    tf = content.text_frame
    tf.text = "Horizontal & Vertical Scaling Strategies"
    
    p = tf.add_paragraph()
    p.text = "Application Layer:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  - Stateless JWT enables horizontal scaling"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "  - Load balancer for multiple instances"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "  - Async processing for non-critical operations"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Database Layer:"
    p.space_before = Pt(12)
    p = tf.add_paragraph()
    p.text = "  - MongoDB sharding for large datasets"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  - Replica sets for high availability"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  - Index optimization for query performance"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Caching Layer:"
    p.space_before = Pt(12)
    p = tf.add_paragraph()
    p.text = "  - Redis for distributed caching (replace Caffeine)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  - Cache warming strategies"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Rate Limiting:"
    p.space_before = Pt(12)
    p = tf.add_paragraph()
    p.text = "  - Distributed rate limiting (Redis-based)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  - API Gateway for centralized rate limiting"
    p.level = 1
    
    # Slide 15: Learning
    slide15 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide15.shapes.title
    title.text = "Key Learnings"
    
    content = slide15.placeholders[1]
    tf = content.text_frame
    tf.text = "Technical & Architectural Insights"
    
    p = tf.add_paragraph()
    p.text = "• Migration from Jakarta EE to Spring Boot ecosystem"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• MongoDB event listeners for automatic audit logging"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Cursor-based pagination vs offset-based pagination"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Multi-layer security: JWT + API Key + CORS"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• PII encryption with key versioning for rotation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Rate limiting strategies for API protection"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Caching strategies for performance optimization"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• OTP-based authentication without password storage"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Update request workflow for data integrity"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• React frontend integration with Spring Boot REST API"
    p.level = 1
    
    return prs

if __name__ == "__main__":
    prs = create_presentation()
    output_file = "Kitchensink_Application_Presentation.pptx"
    prs.save(output_file)
    print(f"Presentation created successfully: {output_file}")

